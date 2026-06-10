"""
material_parser.py — DOCX/PDF/PPTX/TXT → Gemini AI → To'liq dars
=================================================================
Ish jarayoni:
  1. Fayldan matn ajratib olinadi
  2. Gemini mavzular ro'yxatini chiqaradi (yoki regex fallback)
  3. Tanlangan mavzu bo'yicha Gemini to'liq dars yaratadi:
     - title, subject, summary, wow_facts
     - slides (prezentatsiya)
     - quiz (test savollar)
     - flashcards (kartochkalar)
     - key_terms_glossary
     - exam_likely_topics
     - practice_questions
     - main_topics
     - suggested_lesson_plan
  4. Natija DB ga saqlanadi + in-memory/Redis store ga
"""

import os
import json
import re
import threading
import time
import logging
from typing import Optional

# ── fayl kutubxonalari ────────────────────────────────────────
try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import docx as python_docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ── Gemini ────────────────────────────────────────────────────
import google.generativeai as genai

_gemini_key = os.getenv("GEMINI_API_KEY", "")
if _gemini_key:
    genai.configure(api_key=_gemini_key)

logger = logging.getLogger("lectio.material_parser")


def _get_model():
    """Har safar yangi model olish — thread-safe"""
    if not _gemini_key:
        raise RuntimeError("GEMINI_API_KEY sozlanmagan. .env fayliga qo'shing.")
    return genai.GenerativeModel("gemini-1.5-pro")


# ── Redis / In-memory store ───────────────────────────────────
_memory_store: dict = {}
_redis_client = None
_use_redis = False

try:
    import redis as _redis_lib
    _rc = _redis_lib.Redis.from_url(
        os.getenv("REDIS_URL", "redis://localhost:6379/0"),
        decode_responses=True,
        socket_connect_timeout=2,
        socket_timeout=2,
    )
    _rc.ping()
    _redis_client = _rc
    _use_redis = True
    logger.info("[material_parser] Redis: ulandi ✓")
except Exception:
    logger.info("[material_parser] Redis yo'q — xotirada saqlanadi")


def _store_set(key: str, value: str):
    if _use_redis and _redis_client:
        try:
            _redis_client.setex(key, 7200, value)   # 2 soat TTL
            return
        except Exception:
            pass
    _memory_store[key] = value


def _store_get(key: str) -> Optional[str]:
    if _use_redis and _redis_client:
        try:
            return _redis_client.get(key)
        except Exception:
            pass
    return _memory_store.get(key)


# ── Progress ──────────────────────────────────────────────────
def update_progress(material_id: str, stage: str, percent: int, message: str):
    _store_set(
        f"material_progress:{material_id}",
        json.dumps({"stage": stage, "percent": percent, "message": message}),
    )


def get_progress(material_id: str) -> Optional[dict]:
    raw = _store_get(f"material_progress:{material_id}")
    return json.loads(raw) if raw else None


def get_topics(material_id: str) -> Optional[dict]:
    raw = _store_get(f"material_topics:{material_id}")
    return json.loads(raw) if raw else None


def get_lesson_result(material_id: str) -> Optional[dict]:
    raw = _store_get(f"material_lesson_result:{material_id}")
    return json.loads(raw) if raw else None


# ── Matn ajratish ─────────────────────────────────────────────
def extract_text(file_path: str, ext: str) -> str:
    """Fayl turini aniqlab, matni ajratib oladi."""
    text = ""
    ext = ext.lower()

    try:
        if ext == ".pdf":
            if not HAS_PDF:
                raise RuntimeError("PyPDF2 o'rnatilmagan")
            reader = PdfReader(file_path)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
            # Skanerlangan PDF uchun OCR
            if len(text.strip()) < 100:
                try:
                    from pdf2image import convert_from_path
                    import pytesseract
                    for img in convert_from_path(file_path):
                        text += pytesseract.image_to_string(img, lang="uzb+rus+eng") + "\n"
                except Exception as ocr_e:
                    logger.warning(f"OCR xatosi: {ocr_e}")

        elif ext == ".docx":
            if not HAS_DOCX:
                raise RuntimeError("python-docx o'rnatilmagan")
            d = python_docx.Document(file_path)
            # Sarlavha darajalari ham saqlanadi
            for para in d.paragraphs:
                if para.text.strip():
                    prefix = ""
                    if para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "").strip()
                        prefix = "#" * (int(level) if level.isdigit() else 1) + " "
                    text += prefix + para.text.strip() + "\n"
            # Jadvallar ichidagi matnlar
            for table in d.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        text += row_text + "\n"

        elif ext == ".pptx":
            if not HAS_PPTX:
                raise RuntimeError("python-pptx o'rnatilmagan")
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slayd {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text.strip() + "\n"

        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        else:
            raise ValueError(f"Noma'lum fayl turi: {ext}")

    except Exception as e:
        raise RuntimeError(f"Fayl o'qishda xatolik ({ext}): {e}")

    if len(text.strip()) < 20:
        raise RuntimeError(
            "Fayldan matn ajratib bo'lmadi — fayl bo'sh yoki skanerlangan. "
            "Matn asosida saqlangan hujjat yuboring."
        )

    return text.strip()


# ── Mavzular ajratish (AI + fallback) ────────────────────────

# Sarlavha darajalari uchun tartibli patternlar (aniqlash uchun)
_HEADING_PATTERNS = [
    # Level 1 — asosiy bo'limlar
    (1, re.compile(r"^(Mavzu|Bob|Bo['']lim|Qism|Глава|Раздел|Тема|Chapter|Section|Part|ТЕМА|РАЗДЕЛ|ГЛАВА)\s*[\d\s.:\-]*(.{0,100})$", re.IGNORECASE)),
    (1, re.compile(r"^(\d{1,2})[.)]\s{1,4}(?!\d{1,2}[.)])(.{4,90})$")),           # "1. Title" (lekin "1.1." emas)
    (1, re.compile(r"^#{1}\s(.{3,80})$")),                                           # # Title
    # Level 2 — kichik bo'limlar
    (2, re.compile(r"^(\d{1,2}\.\d{1,2})[.):\s]+(.{3,80})$")),                     # "1.1 Title"
    (2, re.compile(r"^#{2}\s(.{3,80})$")),                                           # ## Title
    (2, re.compile(r"^[а-яa-z]?\)\s(.{4,80})$")),                                  # "a) subsection"
    # Level 3 — kichik kichik bo'limlar
    (3, re.compile(r"^(\d{1,2}\.\d{1,2}\.\d{1,2})[.):\s]+(.{3,80})$")),           # "1.1.1 Title"
    (3, re.compile(r"^#{3}\s(.{3,80})$")),                                           # ### Title
]

_CAPS_PAT = re.compile(r"^[А-ЯA-ZO'UG'\s]{4,80}$")


def _build_toc(text: str) -> list:
    """
    Matndan Mundarija (TOC) qurib beradi.
    Har element: {line_idx, level, title, clean_title}
    Bu funksiya hujjat strukturasini aniq tushunadi va
    keyingi `extract_topic_section` uchun pozitsiya ma'lumotini saqlaydi.
    """
    lines = text.splitlines()
    toc = []
    seen = set()

    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped or len(stripped) < 3 or len(stripped) > 140:
            continue

        key = stripped.lower()[:80]
        if key in seen:
            continue

        # Pattern bo'yicha daraja aniqlash
        matched_level = None
        clean_title = stripped

        for level, pat in _HEADING_PATTERNS:
            m = pat.match(stripped)
            if m:
                groups = [g for g in m.groups() if g and len(g.strip()) >= 2]
                clean_title = groups[-1].strip() if groups else stripped
                # Agar sarlavha matni juda qisqa bo'lsa — o'tkazib yubor
                if len(clean_title) < 3:
                    clean_title = stripped
                matched_level = level
                break

        # ALL CAPS qatorlari (1-darajali sarlavha, ko'p hujjatlarda)
        if matched_level is None and _CAPS_PAT.match(stripped):
            word_count = len(stripped.split())
            if 1 <= word_count <= 10:
                matched_level = 1
                clean_title = stripped.title()

        # Ikkala tomondagi bo'sh qatorlar bilan o'ralgan qisqa qatorlar
        if matched_level is None and 5 <= len(stripped) <= 80:
            prev_empty = (i == 0 or not lines[i - 1].strip())
            next_empty = (i >= len(lines) - 1 or not lines[i + 1].strip())
            if prev_empty and next_empty and not stripped.endswith(".") and not stripped.endswith(","):
                word_count = len(stripped.split())
                if 1 <= word_count <= 8:
                    matched_level = 2
                    clean_title = stripped

        if matched_level is not None:
            toc.append({
                "line_idx": i,
                "level": matched_level,
                "title": stripped,
                "clean_title": clean_title,
            })
            seen.add(key)

    return toc


def _regex_extract_headings(text: str) -> list:
    """Matndan sarlavhalar ro'yxatini qaytaradi — AI uchun hint sifatida."""
    toc = _build_toc(text)
    return [item["title"] for item in toc]


def extract_topics_from_text(text: str) -> list:
    """
    Kuchaytirilgan multi-pass strategiya bilan mavzular ro'yxatini chiqaradi:
    1-pass: TOC (hujjat tarkibi, eng aniq)
    2-pass: AI + TOC hint'lari (to'ldiruvchi)
    3-pass: Regex fallback
    """

    # ── 1. TOC asosida ─────────────────────────────────────────
    toc = _build_toc(text)
    toc_topics = [item["title"] for item in toc]
    logger.info(f"[topics] TOC: {len(toc_topics)} sarlavha topildi")

    # Agar TOC yetarlicha kuchli bo'lsa — AI ni to'ldiruvchi sifatida ishlatamiz
    use_ai_supplement = True

    if len(toc_topics) >= 20:
        # Ko'p sarlavha — AI ni faqat tekshirish uchun ishlatamiz
        use_ai_supplement = False
        logger.info(f"[topics] TOC yetarli ({len(toc_topics)} ta), AI skip")
        return toc_topics[:50]

    # ── 2. AI + TOC hint'lari ──────────────────────────────────
    total_len = len(text)
    chunk_size = 12000
    parts = []
    for start in [0, total_len // 4, total_len // 2, 3 * total_len // 4]:
        parts.append(text[start : start + chunk_size])
    combined = "\n---\n".join(dict.fromkeys(parts))[:48000]

    hints_str = ""
    if toc_topics:
        hints_str = f"""
Hujjatdan quyidagi sarlavhalar aniq topildi — BULARNI MAJBURIY kirgiz:
{chr(10).join(f"  - {h}" for h in toc_topics[:30])}

Ular bilan birga matndan boshqa mavzularni ham qo'sh.
"""

    prompt = f"""Sen metodichka va darsliklarni tahlil qiladigan AI mutaxassis.
Quyidagi matn O'zbekiston universiteti metodichkasi.

VAZIFA: Barcha mavzular, boblar, bo'limlar va kichik bo'limlarni to'liq ro'yxat qilib ber.

QOIDALAR:
- Har bir mavzu/bob/bo'lim alohida element — hatto kichik bo'lsa ham
- Sarlavha/nomini AYNAN matnda qanday yozilgan bo'lsa shunday yoz
- Raqamli tartibni SAQLAGAN holda ber (1. 1.1. 2. va h.k.)
- Kamida 5, ko'pi bilan 50 ta element
- FAQAT JSON qaytargin — boshqa hech narsa yo'q

{hints_str}

JSON format:
{{"topics": ["1. Kirish", "1.1 Asosiy tushunchalar", "2. Amaliy qism"]}}

MATN (ko'p qismlar):
{combined}"""

    try:
        model = _get_model()
        resp = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                response_mime_type="application/json",
                max_output_tokens=3000,
                temperature=0.1,
            ),
        )
        raw = resp.text.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        data = json.loads(raw)
        ai_topics = data.get("topics", [])
        if isinstance(ai_topics, list) and len(ai_topics) >= 2:
            cleaned = [str(t).strip() for t in ai_topics if str(t).strip() and 3 <= len(str(t).strip()) <= 120]
            if cleaned:
                # TOC va AI natijalarini birlashtirish (dedup)
                merged = list(dict.fromkeys(toc_topics + cleaned))
                logger.info(f"[topics] AI+TOC birlashtirildi: {len(merged)} mavzu")
                return merged[:50]
    except Exception as e:
        logger.error(f"[topics] Gemini xatosi: {e}")

    # ── 3. TOC fallback ────────────────────────────────────────
    if len(toc_topics) >= 3:
        logger.info(f"[topics] TOC fallback: {len(toc_topics)} mavzu")
        return toc_topics[:50]

    return _fallback_extract_topics(text)


def _fallback_extract_topics(text: str) -> list:
    """Kuchaytirilgan regex asosida mavzular ajratish."""
    lines = text.splitlines()
    topics = []
    seen = set()

    def add(t: str):
        k = t.lower().strip()
        if k not in seen and len(t.strip()) >= 4:
            seen.add(k)
            topics.append(t.strip())

    # Pass 1: Raqamlangan ro'yxat (1. 2) 3- va h.k.)
    num_pattern = re.compile(r"^(\d{1,3}[\d.]*[.):\-]\s*.{4,100})$")
    for line in lines:
        m = num_pattern.match(line.strip())
        if m:
            add(m.group(1))

    # Pass 2: Sarlavha belgilari (# ## ###)
    for line in lines:
        if line.startswith("#") and 4 < len(line.strip()) < 110:
            add(line.lstrip("#").strip())

    # Pass 3: KATTA HARF qatorlari (≥3 so'z)
    for line in lines:
        s = line.strip()
        words = s.split()
        if s.isupper() and 4 < len(s) < 90 and len(words) >= 2:
            add(s.title())

    # Pass 4: "Mavzu", "Bob", "Раздел", "Chapter" kalit so'zlari
    key_pattern = re.compile(
        r"^(Mavzu|Bob|Bo'lim|Qism|Chapter|Раздел|Тема|Part|Section)\s+\d*[.):\s]?.{0,80}$",
        re.IGNORECASE
    )
    for line in lines:
        if key_pattern.match(line.strip()) and 4 < len(line.strip()) < 110:
            add(line.strip())

    # Pass 5: ":" bilan tugagan qatorlar
    for line in lines:
        s = line.strip()
        if s.endswith(":") and 4 < len(s) < 90:
            add(s[:-1])

    if len(topics) >= 3:
        return topics[:50]

    # Pass 6: Bo'sh bo'lmagan qisqa qatorlar (sarlavha bo'lishi mumkin)
    short_lines = [l.strip() for l in lines if 5 < len(l.strip()) < 80 and not l.strip().endswith(".")]
    for l in short_lines[:20]:
        add(l)

    if topics:
        return topics[:50]

    # Oxirgi chora: birinchi 15 ta qator
    first_lines = [l.strip() for l in lines if l.strip() and len(l.strip()) > 4][:15]
    if first_lines:
        return first_lines

    raise RuntimeError(
        "Mavzular topilmadi. Hujjatda sarlavha yoki raqamlangan bo'limlar bo'lishi kerak."
    )


# ── Mavzuga tegishli asl matnni ajratish (TOC asosida) ───────
def _topic_match_score(topic_name: str, entry_title: str) -> float:
    """Mavzu nomi va TOC elementi qanchalik mos kelishini hisoblaydi (0.0–1.0)."""
    # Raqamlarni tozalash
    def _clean(s: str) -> str:
        return re.sub(r"^\d+[\d.]*[.):\-\s]+", "", s).strip().lower()

    tc = _clean(topic_name)
    ec = _clean(entry_title)

    # To'liq mos kelish
    if tc == ec:
        return 1.0

    # Prefix mos kelish (biri boshqasining prefiksi)
    if tc.startswith(ec[:min(15, len(ec))]) or ec.startswith(tc[:min(15, len(tc))]):
        return 0.85

    # So'zlar bo'yicha mos kelish
    tc_words = set(w for w in tc.split() if len(w) >= 4)
    ec_words = set(w for w in ec.split() if len(w) >= 4)
    if tc_words and ec_words:
        overlap = len(tc_words & ec_words)
        score = overlap / max(len(tc_words), len(ec_words))
        return score

    # Bitta qisqa so'z moslik
    if tc and ec and (tc in ec or ec in tc):
        return 0.6

    return 0.0


def extract_topic_section(text: str, topic_name: str) -> str:
    """
    TOC asosida berilgan mavzuga tegishli aniq matn bo'limini ajratadi.
    Strategiya:
    1. TOC quriladi — har bir sarlavhaning qator raqami va darajasi aniqlanadi
    2. Mavzu nomi bo'yicha eng mos TOC elementi topiladi
    3. O'sha darajadagi keyingi sarlavhagacha matn olinadi
    4. Matn 8000 belgiga cheklanadi
    """
    toc = _build_toc(text)
    lines = text.splitlines()

    if not toc:
        # TOC yo'q — oddiy qidiruv
        return _extract_by_keyword(text, topic_name)

    # Eng mos TOC elementini topish
    best_idx = -1
    best_score = 0.0
    for i, entry in enumerate(toc):
        score = _topic_match_score(topic_name, entry["title"])
        if score > best_score:
            best_score = score
            best_idx = i

    if best_idx == -1 or best_score < 0.2:
        return _extract_by_keyword(text, topic_name)

    start_entry = toc[best_idx]
    start_line = start_entry["line_idx"]
    level = start_entry["level"]

    # O'sha yoki yuqori darajadagi keyingi sarlavhani topish
    end_line = len(lines)
    for j in range(best_idx + 1, len(toc)):
        if toc[j]["level"] <= level:
            end_line = toc[j]["line_idx"]
            break

    section = "\n".join(lines[start_line:end_line]).strip()

    if len(section) > 8000:
        section = section[:8000].rsplit("\n", 1)[0] + "\n\n[...matn davom etmoqda...]"

    return section or _extract_by_keyword(text, topic_name)


def _extract_by_keyword(text: str, topic_name: str) -> str:
    """TOC ishlamasa — kalit so'zlar bo'yicha qidiruv (fallback)."""
    clean_topic = re.sub(r"^\d+[\d.]*[.):\-\s]+", "", topic_name).strip()
    topic_words = [w.lower() for w in re.split(r"[\s,;]+", clean_topic) if len(w) >= 4]

    if not topic_words:
        return text[:3000]

    lines = text.splitlines()
    best_score = 0
    start_idx = 0

    for i, line in enumerate(lines):
        line_lower = line.lower()
        score = sum(1 for w in topic_words if w in line_lower)
        if score > best_score:
            best_score = score
            start_idx = i

    end_idx = min(len(lines), start_idx + 200)
    heading_pat = re.compile(r"^(\d{1,3}[.):\-]|#{1,4}\s)")
    for j in range(start_idx + 4, min(start_idx + 200, len(lines))):
        line_j = lines[j].strip()
        if line_j and heading_pat.match(line_j):
            if not any(w in line_j.lower() for w in topic_words[:2]):
                end_idx = j
                break

    section = "\n".join(lines[start_idx:end_idx]).strip()
    if len(section) > 6000:
        section = section[:6000].rsplit("\n", 1)[0] + "\n\n[...matn davomi...]"
    return section or f"[{topic_name} bo'yicha matn topilmadi]"


# ── PPTX → Slaydlar ro'yxati ──────────────────────────────────
def parse_pptx_to_slides(file_path: str) -> list:
    """PPTX prezentatsiyasini slaydlar ro'yxatiga aylantiradi."""
    if not HAS_PPTX:
        raise RuntimeError("python-pptx o'rnatilmagan. `pip install python-pptx` qiling.")

    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        content_parts = []

        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue
            shape_text = shape.text.strip()
            is_title = False
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                if shape.placeholder_format.idx in (0, 1):
                    is_title = True
            if is_title and not title:
                title = shape_text
            elif shape_text != title:
                content_parts.append(shape_text)

        if not title and not content_parts:
            continue

        slides.append({
            "id": i,
            "slide_number": i,
            "title": title or f"Slayd {i}",
            "content": "\n".join(content_parts),
            "notes": "",
            "duration_min": 5,
        })

    return slides


# ── To'liq dars yaratish (AI) ────────────────────────────────
def synthesize_full_lesson(text: str, topic_name: str, professor_id: int) -> dict:
    """
    Berilgan matn va mavzu bo'yicha to'liq dars paketi yaratadi.
    Qaytaruvchi struktura create-lesson sahifasi bilan mos keladi.
    """
    # Mavzuga oid qismni ajratish — ko'proq matn
    total_len = len(text)
    # Mavzuni matnda qidiramiz
    topic_lower = topic_name.lower()
    topic_pos = text.lower().find(topic_lower[:20]) if len(topic_lower) >= 5 else -1
    if topic_pos > 1000:
        # Mavzu atrofidagi matnni olish
        start = max(0, topic_pos - 500)
        chunk = text[start : start + 70000]
    else:
        chunk = text[:70000]

    prompt = f"""Sen professional O'zbekiston universiteti professor-o'qituvchisi va AI dars yaratuvchisan.

VAZIFA: Quyidagi metodichka matni asosida "{topic_name}" mavzusidan TO'LIQ VA BATAFSIL dars paketi yarat.

QOIDALAR:
- Matndan haqiqiy ma'lumot ol, o'ylab chiqarma
- O'ZBEK tilida yaz (termin/atamalar ruscha/inglizcha bo'lishi mumkin)
- Kamida 5-6 ta slayd yarat (har biri 3-5 gap)
- Kamida 5 ta quiz savoli
- Kamida 5 ta flashcard
- Barcha izohlar batafsil bo'lsin

MATN:
{chunk}

TO'LIQ dars paketini QAT'IY JSON formatida yubor. Boshqa hech narsa yozma.

JSON strukturasi (barcha maydonlar to'ldirilishi SHART):
{{
  "title": "{topic_name}",
  "subject": "Fan nomi (masalan: Informatika, Fizika, Matematika...)",
  "level": "Kurs yoki daraja (masalan: Bakalavr 2-kurs)",
  "summary": "Mavzu haqida 3-5 gap qisqacha xulosa. O'zbekiston universitetlari uchun moslashtirilgan.",
  "wow_facts": [
    "Bu mavzu bilan bog'liq hayratlanarli fakt 1",
    "Bu mavzu bilan bog'liq hayratlanarli fakt 2",
    "Bu mavzu bilan bog'liq hayratlanarli fakt 3"
  ],
  "slides": [
    {{
      "id": 1,
      "slide_number": 1,
      "title": "Kirish: {topic_name}",
      "content": "Bugungi mavzu: {topic_name}. Dars maqsadlari va rejasi. Ushbu darsda nimalarga e'tibor beramiz va talabalar nimalarni o'rganadi.",
      "notes": "Talabalarni salomlang va mavzu bilan tanishtiring. Dars maqsadlarini e'lon qiling.",
      "duration_min": 5
    }},
    {{
      "id": 2,
      "slide_number": 2,
      "title": "Nazariy asoslar va ta'riflar",
      "content": "Mavzuning asosiy nazariy tushunchalari, ta'riflar va terminologiya. Muhim atamalar va ularning ma'nosi.",
      "notes": "Ta'riflarni yozma tarzda ko'rsating. Talabalar daftarlariga yozdirsin.",
      "duration_min": 10
    }},
    {{
      "id": 3,
      "slide_number": 3,
      "title": "Asosiy qoidalar va prinsiplar",
      "content": "Mavzuning asosiy qoidalari, prinsiplari va formulalari. Nazariy asoslar chuqurroq ko'rib chiqiladi.",
      "notes": "Formulalar va qoidalarni taxtaga yozing. Misol ko'rsating.",
      "duration_min": 10
    }},
    {{
      "id": 4,
      "slide_number": 4,
      "title": "Amaliy misollar",
      "content": "Mavzuni amalda qo'llash. Konkret misollar, masalalar va ularning yechimlari. Haqiqiy hayotdagi qo'llanish holatlari.",
      "notes": "Talabalar bilan birga misol yeching. Savollar bering.",
      "duration_min": 12
    }},
    {{
      "id": 5,
      "slide_number": 5,
      "title": "Qiyinchiliklar va xatolar",
      "content": "Ko'p uchraydigan xatolar va qiyinchiliklar. Ularni qanday bartaraf etish mumkin. Eslatmalar va ogohlantirishlar.",
      "notes": "Talabalar ko'pincha shu joylarda adashadi — alohida e'tibor bering.",
      "duration_min": 8
    }},
    {{
      "id": 6,
      "slide_number": 6,
      "title": "Xulosa va mustahkamlash",
      "content": "Bugungi dars xulosasi. Asosiy fikrlar va muhim nuqtalar. Keyingi dars mavzusi va uyga vazifa.",
      "notes": "Quiz o'tkazing. Uyga vazifani bering. Savollarga javob bering.",
      "duration_min": 5
    }}
  ],
  "quiz": [
    {{
      "question": "Savol matni (aniq va tushunarli)?",
      "options": ["A) To'g'ri javob", "B) Noto'g'ri variant", "C) Boshqa variant", "D) Noto'g'ri variant 2"],
      "correct": "A",
      "explanation": "To'g'ri javob izohlanishi.",
      "time_limit": 20,
      "points": 1000
    }},
    {{
      "question": "Ikkinchi savol?",
      "options": ["A) Variant 1", "B) To'g'ri javob", "C) Variant 3", "D) Variant 4"],
      "correct": "B",
      "explanation": "Ikkinchi savol izohi.",
      "time_limit": 25,
      "points": 1000
    }},
    {{
      "question": "Uchinchi savol?",
      "options": ["A) Variant", "B) Variant", "C) To'g'ri", "D) Variant"],
      "correct": "C",
      "explanation": "Uchinchi savol izohi.",
      "time_limit": 20,
      "points": 1000
    }},
    {{
      "question": "To'rtinchi savol?",
      "options": ["A) Variant", "B) Variant", "C) Variant", "D) To'g'ri"],
      "correct": "D",
      "explanation": "To'rtinchi savol izohi.",
      "time_limit": 20,
      "points": 1000
    }},
    {{
      "question": "Beshinchi savol?",
      "options": ["A) To'g'ri", "B) Variant", "C) Variant", "D) Variant"],
      "correct": "A",
      "explanation": "Beshinchi savol izohi.",
      "time_limit": 30,
      "points": 1200
    }}
  ],
  "flashcards": [
    {{"front": "Atama yoki tushuncha", "back": "Ta'rif yoki izoh (to'liq)"}},
    {{"front": "Ikkinchi atama", "back": "Ikkinchi ta'rif"}},
    {{"front": "Uchinchi atama", "back": "Uchinchi ta'rif"}},
    {{"front": "To'rtinchi atama", "back": "To'rtinchi ta'rif"}},
    {{"front": "Beshinchi atama", "back": "Beshinchi ta'rif"}}
  ],
  "main_topics": [
    {{
      "title": "Kichik mavzu 1",
      "subtopics": ["Kichik mavzu 1.1", "Kichik mavzu 1.2"],
      "key_concepts": ["Tushuncha 1", "Tushuncha 2"],
      "definitions": [{{"term": "Atama", "definition": "Ta'rif"}}],
      "formulas": ["Formula yoki qoida (agar bo'lsa)"],
      "examples": ["Amaliy misol"]
    }}
  ],
  "suggested_lesson_plan": [
    {{"lesson_number": 1, "title": "1-dars sarlavhasi", "topics": ["mavzu 1", "mavzu 2"], "duration_minutes": 45}},
    {{"lesson_number": 2, "title": "2-dars sarlavhasi", "topics": ["mavzu 3"], "duration_minutes": 45}}
  ],
  "difficulty_distribution": {{"easy": 30, "medium": 50, "hard": 20}},
  "key_terms_glossary": [
    {{"term": "Birinchi atama", "definition": "To'liq ta'rif"}},
    {{"term": "Ikkinchi atama", "definition": "To'liq ta'rif"}},
    {{"term": "Uchinchi atama", "definition": "To'liq ta'rif"}},
    {{"term": "To'rtinchi atama", "definition": "To'liq ta'rif"}},
    {{"term": "Beshinchi atama", "definition": "To'liq ta'rif"}}
  ],
  "exam_likely_topics": [
    "Imtihonda chiqishi mumkin bo'lgan 1-mavzu",
    "Imtihonda chiqishi mumkin bo'lgan 2-mavzu",
    "Imtihonda chiqishi mumkin bo'lgan 3-mavzu",
    "Imtihonda chiqishi mumkin bo'lgan 4-mavzu",
    "Imtihonda chiqishi mumkin bo'lgan 5-mavzu"
  ],
  "practice_questions": [
    {{"question": "Amaliy savol 1?", "answer": "To'liq javob 1", "type": "short_answer"}},
    {{"question": "Amaliy savol 2?", "answer": "To'liq javob 2", "type": "short_answer"}},
    {{"question": "Amaliy savol 3?", "answer": "To'liq javob 3", "type": "short_answer"}}
  ],
  "total_pages": 10
}}

ESLATMALAR:
- Barcha matnlar O'ZBEK tilida bo'lsin
- Matndan olingan haqiqiy ma'lumotlardan foydalanin
- Har bir slide kamida 2-3 gap bo'lsin
- Quiz savollar haqiqiy materialga asoslangan bo'lsin
- Faqat JSON qaytargin, boshqa hech narsa yozma"""

    try:
        model = _get_model()
        resp = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                response_mime_type="application/json",
                max_output_tokens=8000,
                temperature=0.7,
            ),
        )
        raw = resp.text.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        result = json.loads(raw)

        # Minimal validation
        if not isinstance(result, dict):
            raise ValueError("AI noto'g'ri format qaytardi")

        # Ensure required keys exist
        _ensure_lesson_defaults(result, topic_name)
        return result

    except Exception as e:
        logger.error(f"[lesson] Gemini xatosi: {e}")
        return _fallback_full_lesson(text, topic_name)


def _ensure_lesson_defaults(result: dict, topic_name: str):
    """Kamida minimal struktura bo'lishini ta'minlaydi."""
    if not result.get("title"):
        result["title"] = topic_name
    if not result.get("summary"):
        result["summary"] = f"{topic_name} mavzusining qisqacha xulosasi."
    if not result.get("wow_facts"):
        result["wow_facts"] = [f"{topic_name} bo'yicha qiziqarli fakt!"]
    if not result.get("slides"):
        result["slides"] = _default_slides(topic_name)
    if not result.get("quiz"):
        result["quiz"] = _default_quiz(topic_name)
    if not result.get("flashcards"):
        result["flashcards"] = _default_flashcards(topic_name)
    if not result.get("key_terms_glossary"):
        result["key_terms_glossary"] = [{"term": topic_name, "definition": f"{topic_name} — o'quv materialining asosiy tushunchasi."}]
    if not result.get("exam_likely_topics"):
        result["exam_likely_topics"] = [topic_name, f"{topic_name} asosiy tushunchalari"]
    if not result.get("main_topics"):
        result["main_topics"] = [{"title": topic_name, "subtopics": [], "key_concepts": [], "definitions": [], "formulas": [], "examples": []}]


def _default_slides(topic: str) -> list:
    return [
        {"id": 1, "slide_number": 1, "title": f"Kirish: {topic}", "content": f"Bugungi mavzu: {topic}. Dars maqsadlari va rejasi bilan tanishamiz.", "notes": "Talabalarni mavzu bilan tanishtiring.", "duration_min": 5},
        {"id": 2, "slide_number": 2, "title": "Nazariy asoslar", "content": f"{topic} bo'yicha asosiy nazariy bilimlar va ta'riflar. Muhim tushunchalar va formulalar.", "notes": "Ta'riflarga e'tibor bering.", "duration_min": 15},
        {"id": 3, "slide_number": 3, "title": "Amaliy misollar", "content": f"{topic} ni amalda qo'llash. Misollar, masalalar va yechimlar.", "notes": "Talabalar bilan birga yechsin.", "duration_min": 15},
        {"id": 4, "slide_number": 4, "title": "Xulosa va test", "content": f"Bugungi dars xulosasi. {topic} bo'yicha asosiy fikrlar va test savollari.", "notes": "Quiz o'tkazing.", "duration_min": 10},
    ]


def _default_quiz(topic: str) -> list:
    return [
        {"question": f"{topic} nima?", "options": [f"{topic} — bu ta'lim tushunchasi", "Bu boshqa narsa", "Noto'g'ri ta'rif", "Umuman aloqasiz"], "correct": "A", "explanation": f"{topic} asosiy ta'rifi.", "time_limit": 20, "points": 1000},
        {"question": f"{topic} qayerda qo'llaniladi?", "options": ["Hech qayerda", "Faqat laboratoriyada", "Amaliyotda keng", "Faqat nazariyada"], "correct": "C", "explanation": "Amaliyotda keng qo'llaniladi.", "time_limit": 20, "points": 1000},
        {"question": f"{topic} o'rganish nima uchun muhim?", "options": ["Muhim emas", "Imtihon uchun", "Amaliy ko'nikmalar uchun", "Vaqt o'tkazish uchun"], "correct": "C", "explanation": "Amaliy ko'nikmalar kasb rivojida zarur.", "time_limit": 25, "points": 1000},
    ]


def _default_flashcards(topic: str) -> list:
    return [
        {"front": f"{topic} ta'rifi", "back": f"{topic} — bu o'quv dasturining muhim qismi bo'lib, amaliy va nazariy jihatlarga ega."},
        {"front": f"{topic} xususiyatlari", "back": "Tizimlilik, amaliylik, nazariy asoslanganlik."},
        {"front": f"{topic} qo'llanilishi", "back": "Ilm-fan, texnologiya, sanoat va kundalik hayotda."},
    ]


def _fallback_full_lesson(text: str, topic_name: str) -> dict:
    """Gemini ishlamasa ham to'liq struktura qaytaradi."""
    logger.warning(f"[lesson] Fallback ishlatilmoqda: {topic_name}")
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    lower = topic_name.lower()

    related = [l for l in lines if lower in l.lower()][:20]
    summary_parts = related[:4] if related else lines[:4]
    summary = " ".join(summary_parts)
    if len(summary) > 600:
        summary = summary[:600].rsplit(" ", 1)[0] + "..."

    headings = []
    for l in lines:
        if re.match(r"^(\d+[.):\s]|\#{1,3}\s).{5,80}$", l):
            clean = re.sub(r"^(\d+[.):\s]|\#{1,3}\s)", "", l).strip()
            if clean and len(clean) > 4:
                headings.append(clean)
        if len(headings) >= 6:
            break

    if not headings:
        headings = [topic_name, f"{topic_name} asoslari", f"{topic_name} amaliyoti", f"{topic_name} xulosasi"]

    slides = []
    for i, h in enumerate(headings[:5], 1):
        rel_lines = [l for l in lines if h.lower()[:5] in l.lower()][:3]
        content = " ".join(rel_lines) if rel_lines else f"{h} bo'yicha batafsil ma'lumot va misollar."
        if len(content) > 400:
            content = content[:400] + "..."
        slides.append({"id": i, "slide_number": i, "title": h, "content": content, "notes": f"{h} uchun eslatma.", "duration_min": 10})

    quiz = []
    for i, h in enumerate(headings[:5]):
        letters = ["A", "B", "C", "D"]
        opts = [f"{h} — to'g'ri ta'rif", f"Bu boshqa mavzu bilan bog'liq", f"Noto'g'ri ta'rif", f"Umuman aloqasiz narsa"]
        quiz.append({"question": f"{h} nima?", "options": opts, "correct": "A", "explanation": f"{h} — {topic_name} mavzusining muhim qismi.", "time_limit": 20, "points": 1000})

    flashcards = [{"front": h, "back": f"{h} — {topic_name} fanining asosiy tushunchasi."} for h in headings[:6]]

    glossary = [{"term": h, "definition": f"{h}: {topic_name} bilan bog'liq muhim atama."} for h in headings[:5]]

    return {
        "title": topic_name,
        "subject": "Umumiy fan",
        "level": "Bakalavr",
        "summary": summary or f"{topic_name} mavzusi bo'yicha qisqacha tavsif.",
        "wow_facts": [f"{topic_name} bo'yicha qiziqarli fakt: bu soha zamonaviy texnologiyada keng qo'llaniladi!"],
        "slides": slides or _default_slides(topic_name),
        "quiz": quiz or _default_quiz(topic_name),
        "flashcards": flashcards or _default_flashcards(topic_name),
        "main_topics": [{"title": h, "subtopics": [], "key_concepts": [h], "definitions": [{"term": h, "definition": f"{h} ta'rifi."}], "formulas": [], "examples": [f"{h} misol"]} for h in headings[:4]],
        "suggested_lesson_plan": [{"lesson_number": i + 1, "title": h, "topics": [h], "duration_minutes": 45} for i, h in enumerate(headings[:3])],
        "difficulty_distribution": {"easy": 40, "medium": 40, "hard": 20},
        "key_terms_glossary": glossary,
        "exam_likely_topics": headings[:5],
        "practice_questions": [{"question": f"{h} ni tushuntiring", "answer": f"{h} — {topic_name} bilan bog'liq tushuncha.", "type": "short_answer"} for h in headings[:3]],
        "total_pages": len(lines),
    }


# ── Background thread funksiyalari ────────────────────────────
def run_extract_topics(material_id: str, professor_id: int, file_path: str, ext: str):
    """Fayldan mavzular ajratish (background thread)."""
    try:
        update_progress(material_id, "extracting", 10, "Fayl o'qilmoqda...")
        time.sleep(0.3)

        text = extract_text(file_path, ext)
        char_count = len(text)
        update_progress(material_id, "extracting", 30, f"Matn ajratildi ({char_count:,} belgi). AI tahlil boshlanmoqda...")
        time.sleep(0.3)

        # Matnni keyingi bosqich uchun saqlash
        text_path = file_path + ".extracted.txt"
        with open(text_path, "w", encoding="utf-8") as f:
            f.write(text)

        update_progress(material_id, "analyzing", 55, "AI mavzular ro'yxatini aniqlamoqda...")

        topics = extract_topics_from_text(text)

        update_progress(material_id, "analyzing", 90, f"{len(topics)} ta mavzu topildi!")
        time.sleep(0.3)

        _store_set(f"material_topics:{material_id}", json.dumps({"topics": topics, "text_path": text_path}))
        update_progress(material_id, "done", 100, f"✅ {len(topics)} ta mavzu tayyor! Tanlang va dars yarating.")

        # Asl faylni o'chirish (extracted txt qoladi)
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass

    except Exception as e:
        logger.error(f"[extract_topics] Xatolik: {e}")
        update_progress(material_id, "error", 0, f"Xatolik: {str(e)}")


def run_generate_lesson(material_id: str, professor_id: int, topic_name: str, text_path: str):
    """Mavzu bo'yicha to'liq dars yaratish (background thread)."""
    lesson_key = material_id + "_lesson"
    try:
        update_progress(lesson_key, "reading", 15, f"'{topic_name}' — matn o'qilmoqda...")
        time.sleep(0.3)

        if not os.path.exists(text_path):
            raise FileNotFoundError(f"Matn fayli topilmadi: {text_path}")

        with open(text_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        update_progress(lesson_key, "generating", 35, "Gemini AI dars yaratmoqda (30-60 soniya)...")

        final_json = synthesize_full_lesson(text, topic_name, professor_id)

        # Metodichkadan asl mavzu matnini qo'shish (referat uchun)
        final_json["original_topic_text"] = extract_topic_section(text, topic_name)

        update_progress(lesson_key, "saving", 80, "Dars DB va xotiraga saqlanmoqda...")
        time.sleep(0.3)

        # DB ga saqlash (ixtiyoriy — xatolik bo'lsa ham davom etadi)
        _save_to_db(final_json, topic_name, professor_id)

        # Store ga saqlash
        _store_set(
            f"material_lesson_result:{material_id}",
            json.dumps(final_json, ensure_ascii=False),
        )

        slides_count = len(final_json.get("slides", []))
        quiz_count = len(final_json.get("quiz", []))
        flash_count = len(final_json.get("flashcards", []))

        update_progress(
            lesson_key, "done", 100,
            f"✅ Dars tayyor! {slides_count} slayd · {quiz_count} test · {flash_count} karta"
        )

    except Exception as e:
        logger.error(f"[generate_lesson] Xatolik: {e}")
        update_progress(lesson_key, "error", 0, f"Xatolik: {str(e)}")


def _save_to_db(lesson_json: dict, topic_name: str, professor_id: int):
    """DB ga saqlash — xatolik bo'lsa ham dastur to'xtamaydi."""
    try:
        from database import SessionLocal
        from models.lesson import Lesson
        from models.flashcard import FlashCard
        from models.question import Question, QuestionType
        from datetime import datetime, timezone

        db = SessionLocal()
        try:
            title = lesson_json.get("title", topic_name)
            subject = lesson_json.get("subject", "Umumiy")
            wow_facts = lesson_json.get("wow_facts", [])
            wow_fact = wow_facts[0] if wow_facts else ""

            lesson = Lesson(
                title=title,
                topic=subject,
                content=json.dumps(lesson_json, ensure_ascii=False),
                wow_fact=wow_fact,
                professor_id=professor_id,
                created_at=datetime.now(timezone.utc),
            )
            db.add(lesson)
            db.commit()
            db.refresh(lesson)
            lesson_json["lesson_id"] = lesson.id  # Frontend uchun

            # Dars shabloniga tegishli flashcardlar presentation_data JSON da saqlanadi.
            # Bu yerda DB ga yozilmaydi — talabalar /api/chain/flashcards orqali
            # o'z SR kartalarini alohida yaratadi.

            # Quiz savollar
            for q_data in lesson_json.get("quiz", [])[:10]:
                options = q_data.get("options", [])
                correct_letter = q_data.get("correct", "A")
                # Correct letter → actual answer text
                letters = ["A", "B", "C", "D", "E"]
                correct_idx = letters.index(correct_letter) if correct_letter in letters else 0
                correct_text = ""
                if options and correct_idx < len(options):
                    # Remove letter prefix if present: "A) text" → "text"
                    opt = options[correct_idx]
                    correct_text = re.sub(r"^[A-E][).\s]\s*", "", opt).strip() or opt

                q = Question(
                    lesson_id=lesson.id,
                    question=q_data.get("question", "Savol"),
                    type=QuestionType.multiple_choice,
                    options=[re.sub(r"^[A-E][).\s]\s*", "", o).strip() or o for o in options],
                    correct=correct_text or correct_letter,
                    points=q_data.get("points", 100),
                    time_limit=q_data.get("time_limit", 30),
                    explanation=q_data.get("explanation", ""),
                )
                db.add(q)

            db.commit()
            logger.info(f"[DB] Dars saqlandi: lesson_id={lesson.id}")

        except Exception as db_err:
            db.rollback()
            logger.warning(f"[DB] Saqlashda xatolik (non-fatal): {db_err}")
        finally:
            db.close()

    except Exception as import_err:
        logger.warning(f"[DB] Import xatoligi (non-fatal): {import_err}")


# ── Public API ────────────────────────────────────────────────
def start_extract_topics(material_id: str, professor_id: int, file_path: str, ext: str) -> threading.Thread:
    t = threading.Thread(
        target=run_extract_topics,
        args=(material_id, professor_id, file_path, ext),
        daemon=True,
        name=f"extract-{material_id[:8]}",
    )
    t.start()
    return t


def start_generate_lesson(material_id: str, professor_id: int, topic_name: str, text_path: str) -> threading.Thread:
    t = threading.Thread(
        target=run_generate_lesson,
        args=(material_id, professor_id, topic_name, text_path),
        daemon=True,
        name=f"lesson-{material_id[:8]}",
    )
    t.start()
    return t
